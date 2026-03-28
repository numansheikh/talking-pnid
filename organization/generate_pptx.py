"""
Generate presentation.pptx — Talking P&IDs pitch deck
Run: python3 generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
TEAL       = RGBColor(0x0d, 0x94, 0x88)
DARK_BG    = RGBColor(0x0f, 0x17, 0x2a)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0x94, 0xa3, 0xb8)
SLATE      = RGBColor(0x47, 0x55, 0x69)
OFF_WHITE  = RGBColor(0xe2, 0xe8, 0xf0)
NEAR_BLACK = RGBColor(0x1e, 0x29, 0x3b)

# Standard widescreen (13.33 x 7.5 inches)
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_teal_bar(slide, x, y, w, h=Inches(0.04)):
    """Thin teal accent line."""
    add_rect(slide, x, y, w, h, fill_color=TEAL)


# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)

    # Teal left accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    # Eyebrow
    add_text(s, "TALKING P&IDs", Inches(0.5), Inches(1.4), Inches(8), Inches(0.4),
             font_size=13, bold=True, color=TEAL)

    # Main heading
    add_text(s, "Your drawings.", Inches(0.5), Inches(2.0), Inches(9), Inches(1.0),
             font_size=60, bold=True, color=WHITE)
    add_text(s, "Finally answering back.", Inches(0.5), Inches(2.9), Inches(10), Inches(1.0),
             font_size=60, bold=True, color=TEAL)

    # Sub
    add_text(s,
             "AI-powered Q&A for industrial Piping & Instrumentation Diagrams",
             Inches(0.5), Inches(4.2), Inches(8), Inches(0.6),
             font_size=18, color=LIGHT_GREY)

    # ZeerakLabs bottom right
    add_text(s, "ZeerakLabs  ·  talkingpnid.zeeraklabs.online",
             Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             font_size=11, color=RGBColor(0x47,0x55,0x69), align=PP_ALIGN.RIGHT)

    return s


# ── SLIDE 2: THE PROBLEM QUOTE ────────────────────────────────────────────────
def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "THE PROBLEM", Inches(0.5), Inches(0.6), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)

    lines = [
        ('"Which isolation valves do I need to close?', 28, WHITE, False),
        ('What\'s the design pressure of this vessel?',  28, WHITE, False),
        ('What happens if this valve fails?"',            28, WHITE, False),
        ('',                                              14, WHITE, False),
        ('These questions should take 30 seconds.',       28, TEAL,  True),
        ('Today they take 30 minutes.',                   28, LIGHT_GREY, False),
    ]

    y = Inches(1.3)
    for line, size, color, bold in lines:
        add_text(s, line, Inches(0.8), y, Inches(11), Inches(0.6),
                 font_size=size, bold=bold, color=color, italic=(size==28 and not bold))
        y += Inches(0.52 if size == 28 else 0.25)

    add_text(s, "— Every process facility, every shift",
             Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
             font_size=13, color=SLATE, italic=True)
    return s


# ── SLIDE 3: PAIN POINTS ──────────────────────────────────────────────────────
def slide_pain(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "THE CHALLENGE", Inches(0.5), Inches(0.4), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)
    add_text(s, "P&IDs hold the knowledge. Accessing it is the problem.",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.6),
             font_size=28, bold=True, color=WHITE)

    cards = [
        ("Senior engineers = human search engines",
         "Experienced staff are constantly interrupted for questions that should be self-service."),
        ("Shutdown prep takes days",
         "Isolation lists, instrument checks, and spec break summaries are built manually."),
        ("New engineers take months",
         "Learning complex P&IDs requires sustained mentorship that's hard to scale."),
        ("Field decisions without full information",
         "Engineers on the plant floor often work from memory — increasing error risk."),
    ]

    positions = [
        (Inches(0.5),  Inches(2.1)),
        (Inches(6.8),  Inches(2.1)),
        (Inches(0.5),  Inches(4.5)),
        (Inches(6.8),  Inches(4.5)),
    ]

    for (title, body), (x, y) in zip(cards, positions):
        add_rect(s, x, y, Inches(6.0), Inches(2.1),
                 fill_color=RGBColor(0x1e,0x29,0x3b),
                 line_color=RGBColor(0x33,0x41,0x55), line_width=Pt(0.5))
        add_text(s, title, x+Inches(0.2), y+Inches(0.15), Inches(5.6), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)
        add_text(s, body, x+Inches(0.2), y+Inches(0.65), Inches(5.6), Inches(1.2),
                 font_size=12, color=LIGHT_GREY)
    return s


# ── SLIDE 4: BEFORE / AFTER ───────────────────────────────────────────────────
def slide_before_after(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "THE SHIFT", Inches(0.5), Inches(0.4), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)
    add_text(s, "From hunting drawings to instant answers",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.55),
             font_size=28, bold=True, color=WHITE)

    # Before box
    add_rect(s, Inches(0.5), Inches(1.8), Inches(5.6), Inches(5.0),
             fill_color=RGBColor(0x1f,0x10,0x10),
             line_color=RGBColor(0x7f,0x1d,0x1d), line_width=Pt(0.75))
    add_text(s, "TODAY", Inches(0.7), Inches(1.95), Inches(5), Inches(0.35),
             font_size=11, bold=True, color=RGBColor(0xf8,0x71,0x71))

    befores = [
        "Search through paper or PDF drawings",
        "Manually trace lines and connections",
        "Call a senior engineer for help",
        "Wait — and interrupt their work",
        "Get an answer that may be from memory",
    ]
    for i, line in enumerate(befores):
        add_text(s, "—  " + line, Inches(0.7), Inches(2.4 + i*0.55), Inches(5.2), Inches(0.5),
                 font_size=13, color=RGBColor(0xcd,0xd5,0xe1))

    # Arrow
    add_text(s, "→", Inches(6.3), Inches(4.0), Inches(0.8), Inches(0.6),
             font_size=32, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # After box
    add_rect(s, Inches(7.2), Inches(1.8), Inches(5.6), Inches(5.0),
             fill_color=RGBColor(0x06,0x1e,0x1c),
             line_color=TEAL, line_width=Pt(0.75))
    add_text(s, "WITH TALKING P&IDs", Inches(7.4), Inches(1.95), Inches(5), Inches(0.35),
             font_size=11, bold=True, color=TEAL)

    afters = [
        "Open a browser on any device",
        "Ask the question in plain English",
        "Get an answer with cited tag numbers",
        "Ask follow-up questions naturally",
        "Available 24/7 to every engineer",
    ]
    for i, line in enumerate(afters):
        add_text(s, "→  " + line, Inches(7.4), Inches(2.4 + i*0.55), Inches(5.2), Inches(0.5),
                 font_size=13, color=RGBColor(0xcd,0xd5,0xe1))
    return s


# ── SLIDE 5: DEMO QUESTIONS ───────────────────────────────────────────────────
def slide_demo(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "IN ACTION", Inches(0.5), Inches(0.4), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)
    add_text(s, "Real questions. Real answers.",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.55),
             font_size=28, bold=True, color=WHITE)

    qs = [
        ("What are all the isolation valves for vessel 362-V001?",                    "Maintenance isolation"),
        ("EZV-0002 closes unexpectedly — what are the upstream effects and safeguards?","Troubleshooting"),
        ("List all locked-open and locked-closed valves on this P&ID",                 "Valve lineup check"),
        ("What instruments have high-high or low-low alarm functions?",                "Safety review"),
        ("Show me all spec breaks and their boundary equipment",                        "Integrity check"),
    ]

    for i, (q, tag) in enumerate(qs):
        y = Inches(1.7 + i * 1.02)
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.82),
                 fill_color=RGBColor(0x1a,0x23,0x35),
                 line_color=TEAL, line_width=Pt(1.5))
        # Left teal bar inside card
        add_rect(s, Inches(0.5), y, Inches(0.06), Inches(0.82), fill_color=TEAL)
        add_text(s, f'"{q}"', Inches(0.7), y+Inches(0.12), Inches(10.2), Inches(0.6),
                 font_size=13, color=OFF_WHITE, italic=True)
        add_rect(s, Inches(11.0), y+Inches(0.2), Inches(1.65), Inches(0.42),
                 fill_color=RGBColor(0x0d,0x40,0x3c))
        add_text(s, tag, Inches(11.0), y+Inches(0.2), Inches(1.65), Inches(0.42),
                 font_size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    return s


# ── SLIDE 6: HOW IT WORKS ─────────────────────────────────────────────────────
def slide_how(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "IMPLEMENTATION", Inches(0.5), Inches(0.4), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)
    add_text(s, "From your PDFs to answers in days",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.55),
             font_size=28, bold=True, color=WHITE)

    steps = [
        ("1", "Send us\nyour P&IDs",       "Scanned PDFs in any format — no CAD files needed"),
        ("2", "We process\nthem",           "AI extracts full topology, attributes and relationships"),
        ("3", "Validate with\nyour team",   "Structured review with your engineers to confirm accuracy"),
        ("4", "Your team\nqueries",         "Private, secure web app — any browser, any device"),
    ]

    for i, (num, title, body) in enumerate(steps):
        x = Inches(0.5 + i * 3.2)
        cy = Inches(2.5)
        # Circle
        add_rect(s, x+Inches(0.9), cy, Inches(1.1), Inches(1.1),
                 fill_color=RGBColor(0x0d,0x40,0x3c),
                 line_color=TEAL, line_width=Pt(1.5))
        add_text(s, num, x+Inches(0.9), cy, Inches(1.1), Inches(1.1),
                 font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text(s, title, x, Inches(3.9), Inches(3.0), Inches(0.8),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, body, x, Inches(4.8), Inches(3.0), Inches(1.5),
                 font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

        # Connector arrow (except after last)
        if i < 3:
            add_text(s, "→", x+Inches(2.95), Inches(2.85), Inches(0.5), Inches(0.5),
                     font_size=20, color=TEAL, align=PP_ALIGN.CENTER)
    return s


# ── SLIDE 7: PROOF ────────────────────────────────────────────────────────────
def slide_proof(prs):
    s = blank_slide(prs)
    fill_bg(s, DARK_BG)
    add_rect(s, Inches(0), Inches(0), Inches(0.06), H, fill_color=TEAL)

    add_text(s, "PROOF OF CONCEPT", Inches(0.5), Inches(0.4), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=TEAL)
    add_text(s, "Validated on real production oil field P&IDs",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.55),
             font_size=28, bold=True, color=WHITE)

    metrics = [
        ("3",    "Interconnected P&IDs\nfrom Rumaila oil field, Iraq"),
        ("79%",  "Benchmark accuracy\non 10 engineering questions"),
        ("385",  "Components modelled\nacross all three drawings"),
    ]

    for i, (num, label) in enumerate(metrics):
        x = Inches(0.8 + i * 4.2)
        add_rect(s, x, Inches(1.9), Inches(3.8), Inches(2.4),
                 fill_color=RGBColor(0x06,0x1e,0x1c),
                 line_color=TEAL, line_width=Pt(0.75))
        add_text(s, num, x, Inches(2.0), Inches(3.8), Inches(1.0),
                 font_size=52, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_text(s, label, x, Inches(3.1), Inches(3.8), Inches(1.0),
                 font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    add_text(s,
             "Benchmarked against 10 representative engineering questions covering isolation valve lists, "
             "design conditions, troubleshooting scenarios, alarm functions, locked valve positions, "
             "spec breaks, and cross-diagram connections.",
             Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.5),
             font_size=13, color=SLATE)
    return s


# ── SLIDE 8: PILOT / CTA ──────────────────────────────────────────────────────
def slide_cta(prs):
    s = blank_slide(prs)
    fill_bg(s, TEAL)

    add_text(s, "NEXT STEPS", Inches(0.7), Inches(0.7), Inches(8), Inches(0.35),
             font_size=11, bold=True, color=RGBColor(0xff,0xff,0xff))

    add_text(s, "Ready to try it\non your drawings?",
             Inches(0.7), Inches(1.2), Inches(7.5), Inches(2.2),
             font_size=52, bold=True, color=WHITE)

    bullets = [
        ("48-hour trial",     "Query 3 real P&IDs right now at talkingpnid.zeeraklabs.online"),
        ("Structured pilot",  "We process your P&IDs and work with your team to validate"),
        ("No commitment",     "The pilot gives a clear picture of fit before any contract"),
    ]

    for i, (title, body) in enumerate(bullets):
        y = Inches(3.6 + i * 0.9)
        add_rect(s, Inches(0.7), y, Inches(7.5), Inches(0.75),
                 fill_color=RGBColor(0x0b,0x7a,0x70))
        add_text(s, f"✓  {title} — {body}", Inches(0.9), y+Inches(0.1), Inches(7.2), Inches(0.55),
                 font_size=13, color=WHITE, bold=False)

    # URL box
    add_rect(s, Inches(9.2), Inches(2.0), Inches(3.7), Inches(3.0),
             fill_color=RGBColor(0x0b,0x7a,0x70),
             line_color=WHITE, line_width=Pt(1.0))
    add_text(s, "talkingpnid\n.zeeraklabs.online",
             Inches(9.2), Inches(2.2), Inches(3.7), Inches(1.6),
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "hello@zeeraklabs.com",
             Inches(9.2), Inches(3.9), Inches(3.7), Inches(0.5),
             font_size=13, color=RGBColor(0xcc,0xf2,0xee), align=PP_ALIGN.CENTER)

    add_text(s, "ZeerakLabs  ·  zeeraklabs.com",
             Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
             font_size=11, color=RGBColor(0xff,0xff,0xff), align=PP_ALIGN.RIGHT)
    return s


# ── MAIN ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_pain(prs)
    slide_before_after(prs)
    slide_demo(prs)
    slide_how(prs)
    slide_proof(prs)
    slide_cta(prs)

    out = "presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({prs.slides.__len__()} slides)")
